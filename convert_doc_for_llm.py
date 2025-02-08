# convert_doc_for_llm-v3.py
import os
import mimetypes
import subprocess
from PIL import Image
import pytesseract
import aspose.words as aw  # Note: Aspose.Words might require a license for commercial use.
from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
import pandas as pd
import fitz  # PyMuPDF for PDF handling

def detect_file_type(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    print(f"MimeTypes Guess: {mime_type}") # Debug print
    if not mime_type: # if mime_type is None or '', try to guess based on extension
        _, ext = os.path.splitext(file_path)
        ext_lower = ext.lower()
        print(f"File Extension: {ext_lower}") # Debug print
        if ext_lower in ['.txt', '.log', '.csv']:
            mime_type = 'text/plain'
        elif ext_lower in ['.epub']:
            print("DEBUG: EPUB extension detected!") # Added debug print here
            mime_type = 'application/epub+zip'
        elif ext_lower in ['.pdf']:
            print("DEBUG: PDF extension detected!") # Added debug print here
            mime_type = 'application/pdf'
        elif ext_lower in ['.docx', '.doc']:
            mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif ext_lower in ['.pptx', '.ppt']:
            mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif ext_lower in ['.xlsx', '.xls']:
            mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        elif ext_lower in ['.jpg', '.jpeg', '.png', '.gif', '.tiff', '.bmp']:
            mime_type = 'image/' + ext_lower[1:]
    print(f"Detected File Type (before return): {mime_type}") # Debug print
    return mime_type

def convert_pptx_to_text(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"[SLIDE {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except Exception as e:
        return f"Error converting PPTX: {e}"
    return "\n".join(text)

def convert_docx_to_text(file_path):
    try:
        doc = aw.Document(file_path)
        return doc.get_text()
    except Exception as e:
        return f"Error converting DOCX: {e}"

def convert_xlsx_to_text(file_path):
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet_name, sheet_data in df.items():
            text.append(f"[SHEET: {sheet_name}]")
            text.append(sheet_data.to_string(index=False))
        return "\n\n".join(text)
    except Exception as e:
        return f"Error converting XLSX: {e}"

def convert_epub_to_text(file_path):
    text = []
    try:
        book = epub.read_epub(file_path)
        for item in book.get_items():
            if item.get_type() == 9: # Changed from epub.ITEM_DOCUMENT to 9
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return "\n\n".join(text)
    except Exception as e:
        return f"Error converting EPUB: {e}"

def convert_image_to_text(file_path):
    try:
        return pytesseract.image_to_string(Image.open(file_path))
    except Exception as e:
        return f"Error converting Image (OCR): {e}"

def convert_pdf_to_text(file_path):
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        return f"Error converting PDF: {e}"

def convert_text_file_to_text(file_path): # Directly read text files
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Error reading Text file: {e}"


def convert_to_text(input_file, output_file):
    file_type = detect_file_type(input_file)
    print(f"Detected file type: {file_type}") # User feedback

    if file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = convert_pptx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text = convert_docx_to_text(input_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        text = convert_xlsx_to_text(input_file)
    elif file_type == "application/epub+zip":
        text = convert_epub_to_text(input_file)
    elif file_type.startswith("image/"):
        text = convert_image_to_text(input_file)
    elif file_type == "application/pdf":
        text = convert_pdf_to_text(input_file)
    elif file_type == 'text/plain': # Handle plain text files directly
        text = convert_text_file_to_text(input_file)
    elif file_type is None: # Handle completely unknown types gracefully
        return f"Error: Could not detect file type for: {input_file}. Please ensure it is a supported format (PPTX, DOCX, XLSX, EPUB, Image, PDF, or Text)."
    else:
        return f"Error: Unsupported file type: {file_type}. Please provide a PPTX, DOCX, XLSX, EPUB, Image, PDF, or Text file." # More informative error

    if isinstance(text, str) and text.startswith("Error"): # Check if conversion function returned an error message
        print(text) # Print the specific error message to the user
        return # Exit if there was an error

    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(text)
        print(f"🎉 Conversion complete! Text saved to: {output_file}") # Positive feedback
    except Exception as e:
        print(f"Error writing to output file: {e}")

if __name__ == "__main__": # Make script directly executable
    input_file = input("Enter the path to your input document: ")
    output_file = input("Enter the path to save the output text file: ")
    result = convert_to_text(input_file, output_file) # Capture return value
    if result and result.startswith("Error"): # Check for error message
        print(result) # Print error message

